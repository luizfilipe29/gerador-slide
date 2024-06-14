from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide(presentation, slide_text, font_size, font_color, bg_color):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.add_paragraph()
    paragraph.text = slide_text
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = RGBColor(*map(int, font_color))
    paragraph.font.name = 'Arial'
    paragraph.alignment = PP_ALIGN.LEFT
    
    if bg_color.startswith("#"):
        bg_color = hex_to_rgb(bg_color)
    
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def split_text_into_slides(lyrics, max_lines_per_slide):
    lines = lyrics.split('\n')
    slides_text = []
    current_slide_text = []
    for line in lines:
        if len(current_slide_text) < max_lines_per_slide:
            current_slide_text.append(line)
        else:
            slides_text.append('\n'.join(current_slide_text))
            current_slide_text = [line]
    if current_slide_text:
        slides_text.append('\n'.join(current_slide_text))
    return slides_text

def generate_slides_from_lyrics(lyrics, filename, font_size, font_color, bg_color, max_lines_per_slide):
    presentation = Presentation()
    slides_text = split_text_into_slides(lyrics, max_lines_per_slide)
    for slide_text in slides_text:
        create_slide(presentation, slide_text, font_size, font_color, bg_color)
    presentation.save(f'{filename}.pptx')
